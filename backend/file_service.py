"""File processing service for Sertex.

Handles text/content extraction from uploaded files across formats:
- Documents: PDF, DOCX, TXT, MD
- Spreadsheets: XLSX, CSV
- Presentations: PPTX
- Images: JPG, JPEG, PNG, WEBP, GIF (via GPT-5.2 Vision)
- Audio/Video: WAV, MP3, M4A, WEBM, MP4 (via Whisper)

Also provides summarization and Q&A over the extracted content using the
Emergent Universal LLM Key (OpenAI GPT-5.2).
"""
import os
import io
import base64
import logging
import uuid
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any
from dotenv import load_dotenv
from pydantic import BaseModel, Field

from emergentintegrations.llm.chat import LlmChat, UserMessage, ImageContent

load_dotenv()
logger = logging.getLogger(__name__)


# ---- Constants -----------------------------------------------------------
MAX_UPLOAD_BYTES = 50 * 1024 * 1024  # 50 MB
MAX_TEXT_CHARS = 200_000  # ~50k tokens — safe upper bound for LLM injection
SUMMARY_MAX_CHARS = 60_000  # chars fed to summarizer

DOCUMENT_EXTS = {"pdf", "docx", "doc", "txt", "md", "rtf"}
SPREADSHEET_EXTS = {"xlsx", "xls", "csv"}
PRESENTATION_EXTS = {"pptx", "ppt"}
IMAGE_EXTS = {"jpg", "jpeg", "png", "webp", "gif", "bmp"}
AUDIO_EXTS = {"wav", "mp3", "m4a", "webm", "ogg", "mp4", "mpeg", "mpga"}


def get_category(ext: str) -> str:
    ext = (ext or "").lower()
    if ext in DOCUMENT_EXTS:
        return "document"
    if ext in SPREADSHEET_EXTS:
        return "spreadsheet"
    if ext in PRESENTATION_EXTS:
        return "presentation"
    if ext in IMAGE_EXTS:
        return "image"
    if ext in AUDIO_EXTS:
        return "audio"
    return "other"


def get_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_supported_extension(ext: str) -> bool:
    ext = (ext or "").lower()
    return (
        ext in DOCUMENT_EXTS
        or ext in SPREADSHEET_EXTS
        or ext in PRESENTATION_EXTS
        or ext in IMAGE_EXTS
        or ext in AUDIO_EXTS
    )


# ---- Models --------------------------------------------------------------
class FileRecord(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    user_id: str
    original_filename: str
    storage_path: str
    content_type: str
    extension: str
    category: str  # document | spreadsheet | presentation | image | audio | other
    size: int
    extracted_text: str = ""
    extracted_chars: int = 0
    extraction_status: str = "pending"  # pending | ok | partial | failed | unsupported
    extraction_error: Optional[str] = None
    summary: Optional[str] = None
    is_deleted: bool = False
    rag_status: str = "pending"  # pending | indexing | ok | empty | failed
    rag_chunks: int = 0
    rag_error: Optional[str] = None
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())


# ---- Text extraction -----------------------------------------------------
def _truncate(text: str, limit: int = MAX_TEXT_CHARS) -> str:
    if not text:
        return ""
    if len(text) <= limit:
        return text
    return text[:limit] + "\n\n[... truncated ...]"


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            logger.debug(f"PDF page {i} extract failed: {e}")
            txt = ""
        if txt.strip():
            pages.append(f"--- Sayfa {i} ---\n{txt.strip()}")
    return _truncate("\n\n".join(pages))


def extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts))


def extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sayfa: {sheet_name} ---")
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            values = [
                str(v) if v is not None else "" for v in row
            ]
            if any(values):
                parts.append(" | ".join(values))
                row_count += 1
            if row_count >= 2000:  # avoid runaway
                parts.append(f"[... {sheet_name}: ilk 2000 satır gösteriliyor ...]")
                break
    return _truncate("\n".join(parts))


def extract_csv(data: bytes) -> str:
    import csv
    text = None
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1254"):
        try:
            text = data.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    parts: List[str] = []
    for i, row in enumerate(reader):
        parts.append(" | ".join(row))
        if i >= 5000:
            parts.append("[... ilk 5000 satır gösteriliyor ...]")
            break
    return _truncate("\n".join(parts))


def extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slayt {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return _truncate("\n".join(parts))


def extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1254"):
        try:
            return _truncate(data.decode(enc))
        except UnicodeDecodeError:
            continue
    return _truncate(data.decode("utf-8", errors="replace"))


async def extract_image_via_vision(data: bytes, mime_type: str = "image/jpeg") -> str:
    """Analyze image content using GPT-5.2 vision. Returns Turkish description."""
    key = os.environ.get("EMERGENT_LLM_KEY", "")
    if not key:
        raise RuntimeError("EMERGENT_LLM_KEY not configured")

    b64 = base64.b64encode(data).decode("utf-8")
    session_id = f"vision-{uuid.uuid4()}"
    chat = LlmChat(
        api_key=key,
        session_id=session_id,
        system_message=(
            "Sen görsel analiz uzmanısın. Kullanıcının yüklediği görseli detaylı "
            "biçimde Türkçe olarak analiz et. Şunları belirt: (1) genel içerik, "
            "(2) görülen nesneler ve kişiler, (3) varsa metin (OCR), "
            "(4) renk ve kompozisyon notları, (5) olası bağlam veya niyet. "
            "Yanıtı düz metin, madde işaretsiz ver."
        ),
    ).with_model("openai", "gpt-5.2")

    msg = UserMessage(
        text="Bu görseli detaylı analiz et ve Türkçe yanıt ver. Görselde metin varsa OCR ile aktar.",
        file_contents=[ImageContent(image_base64=b64)],
    )
    response = await chat.send_message(msg)
    return str(response).strip()


async def extract_audio_via_whisper(
    data: bytes, filename: str, content_type: str
) -> str:
    """Transcribe audio via existing whisper_service."""
    from whisper_service import transcribe_audio
    return await transcribe_audio(
        data,
        filename=filename,
        content_type=content_type,
        language="tr",
    )


async def extract_content(
    data: bytes,
    filename: str,
    content_type: str,
) -> Dict[str, Any]:
    """Main entry point. Returns dict with keys:
    { text, status, error, category }
    """
    ext = get_extension(filename)
    category = get_category(ext)
    try:
        if ext == "pdf":
            text = extract_pdf(data)
        elif ext in ("docx",):
            text = extract_docx(data)
        elif ext in ("xlsx", "xls"):
            text = extract_xlsx(data)
        elif ext == "csv":
            text = extract_csv(data)
        elif ext in ("pptx",):
            text = extract_pptx(data)
        elif ext in ("txt", "md", "rtf"):
            text = extract_txt(data)
        elif ext in IMAGE_EXTS:
            mime = content_type if content_type and content_type.startswith("image/") else f"image/{ext}"
            text = await extract_image_via_vision(data, mime_type=mime)
        elif ext in AUDIO_EXTS:
            text = await extract_audio_via_whisper(data, filename, content_type)
        else:
            return {
                "text": "",
                "status": "unsupported",
                "error": f".{ext} formatı henüz desteklenmiyor",
                "category": category,
            }
        return {
            "text": text or "",
            "status": "ok" if (text and text.strip()) else "partial",
            "error": None if (text and text.strip()) else "İçerik boş veya çıkarılamadı",
            "category": category,
        }
    except Exception as e:
        logger.exception(f"Extraction failed for {filename}")
        return {
            "text": "",
            "status": "failed",
            "error": str(e)[:300],
            "category": category,
        }


# ---- LLM helpers ---------------------------------------------------------
def _llm_client(session_id: str, system: str) -> LlmChat:
    key = os.environ.get("EMERGENT_LLM_KEY", "")
    if not key:
        raise RuntimeError("EMERGENT_LLM_KEY not configured")
    return LlmChat(
        api_key=key,
        session_id=session_id,
        system_message=system,
    ).with_model("openai", "gpt-5.2")


async def summarize_text(text: str, filename: str = "") -> str:
    """Produce a concise Turkish summary of an extracted document."""
    if not text or not text.strip():
        return ""
    payload = text[:SUMMARY_MAX_CHARS]
    system = (
        "Sen bir Türkçe belge özetleme uzmanısın. Verilen içerikten kısa, öz ve "
        "aksiyona dönüştürülebilir bir özet çıkar. Şu formatı kullan:\n"
        "**Kısa Özet:** 1-2 cümle.\n"
        "**Ana Bulgular:** 3-6 madde.\n"
        "**Öneriler / Aksiyonlar:** varsa 1-3 madde.\n"
        "Emoji kullanma. Uydurma yapma; içerikte yoksa 'belirtilmemiş' de."
    )
    chat = _llm_client(f"summarize-{uuid.uuid4()}", system)
    resp = await chat.send_message(
        UserMessage(text=f"Dosya: {filename}\n\nİçerik:\n{payload}")
    )
    return str(resp).strip()


async def answer_question(text: str, question: str, filename: str = "") -> str:
    """Answer a user question grounded in the extracted file content."""
    if not text or not text.strip():
        return "Bu dosyanın içeriği çıkarılamamış görünüyor efendim. Farklı bir dosya deneyebilir misiniz?"
    payload = text[:SUMMARY_MAX_CHARS]
    system = (
        "Sen Sertex — kullanıcının kişisel yapay zeka asistanısın. Aşağıdaki dosya "
        "içeriğine dayanarak Türkçe, kısa ve kesin cevaplar ver. Yalnızca "
        "içerikte olan bilgiye dayan; içerikte yoksa 'dosyada belirtilmemiş' de. "
        "Emoji kullanma."
    )
    chat = _llm_client(f"file-qa-{uuid.uuid4()}", system)
    resp = await chat.send_message(
        UserMessage(
            text=(
                f"[Dosya: {filename}]\n[İçerik]:\n{payload}\n\n"
                f"[Kullanıcının sorusu]:\n{question}"
            )
        )
    )
    return str(resp).strip()
